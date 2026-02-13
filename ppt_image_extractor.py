"""
PowerPoint Image Extraction Service
Extracts images from PPT/PPTX files
"""

import os
from pptx import Presentation
from PIL import Image
import io


def extract_images_from_ppt(input_path, output_dir):
    """
    Extract all images from PowerPoint presentation
    
    Args:
        input_path: Path to input PPTX
        output_dir: Directory to save extracted images
        
    Returns:
        dict: Result information
    """
    try:
        # Create output directory if it doesn't exist
        os.makedirs(output_dir, exist_ok=True)
        
        # Load presentation
        prs = Presentation(input_path)
        
        image_count = 0
        extracted_images = []
        
        # Iterate through slides
        for slide_num, slide in enumerate(prs.slides, 1):
            # Iterate through shapes in slide
            for shape_num, shape in enumerate(slide.shapes, 1):
                # Check if shape contains an image
                if hasattr(shape, "image"):
                    image = shape.image
                    
                    # Get image bytes
                    image_bytes = image.blob
                    
                    # Determine format
                    ext = image.ext
                    
                    # Create output filename
                    output_filename = f"slide_{slide_num}_img_{shape_num}.{ext}"
                    output_path = os.path.join(output_dir, output_filename)
                    
                    # Save image
                    with open(output_path, "wb") as img_file:
                        img_file.write(image_bytes)
                    
                    extracted_images.append({
                        'filename': output_filename,
                        'slide': slide_num,
                        'size': len(image_bytes)
                    })
                    
                    image_count += 1
        
        return {
            'success': True,
            'image_count': image_count,
            'extracted_images': extracted_images,
            'output_dir': output_dir
        }
        
    except Exception as e:
        return {
            'success': False,
            'error': str(e)
        }


def extract_images_from_ppt_batch(input_paths, output_base_dir):
    """
    Extract images from multiple PowerPoint files
    
    Args:
        input_paths: List of input PPTX paths
        output_base_dir: Base directory for outputs
        
    Returns:
        list: Results for each presentation
    """
    results = []
    
    for input_path in input_paths:
        filename = os.path.basename(input_path)
        name, ext = os.path.splitext(filename)
        
        # Create subdirectory for this presentation's images
        output_dir = os.path.join(output_base_dir, f"{name}_images")
        
        result = extract_images_from_ppt(input_path, output_dir)
        result['ppt_filename'] = filename
        results.append(result)
    
    return results