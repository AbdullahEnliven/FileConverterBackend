"""
PowerPoint Slide Export Service
Exports slides as images (PNG/JPG)
"""

import os
from pptx import Presentation
from pptx.util import Inches
from PIL import Image, ImageDraw, ImageFont
import io


def export_slides_as_images(input_path, output_dir, image_format='png', slide_numbers=None):
    """
    Export PowerPoint slides as images
    
    Args:
        input_path: Path to input PPTX
        output_dir: Directory to save slide images
        image_format: Output format ('png' or 'jpg')
        slide_numbers: List of slide numbers to export (None = all slides)
        
    Returns:
        dict: Result information
    """
    try:
        # Create output directory if it doesn't exist
        os.makedirs(output_dir, exist_ok=True)
        
        # Load presentation
        prs = Presentation(input_path)
        
        exported_slides = []
        total_slides = len(prs.slides)
        
        # Determine which slides to export
        if slide_numbers is None:
            slides_to_export = range(1, total_slides + 1)
        else:
            slides_to_export = slide_numbers
        
        # Create images for each slide
        for slide_num in slides_to_export:
            if slide_num < 1 or slide_num > total_slides:
                continue
            
            slide_idx = slide_num - 1
            slide = prs.slides[slide_idx]
            
            # Create a white canvas (standard 16:9 aspect ratio)
            width = 1920
            height = 1080
            img = Image.new('RGB', (width, height), color='white')
            draw = ImageDraw.Draw(img)
            
            # Try to use a font
            try:
                font_title = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 48)
                font_text = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 32)
            except:
                font_title = ImageFont.load_default()
                font_text = ImageFont.load_default()
            
            # Draw slide content (simplified - just text)
            y_position = 100
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    
                    # Check if it's a title
                    is_title = hasattr(shape, "name") and "Title" in shape.name
                    font = font_title if is_title else font_text
                    
                    # Draw text (wrap if too long)
                    lines = text.split('\n')
                    for line in lines:
                        if y_position < height - 100:
                            draw.text((100, y_position), line[:80], fill='black', font=font)
                            y_position += 60 if is_title else 40
            
            # Draw slide number
            draw.text((width - 150, height - 80), f"Slide {slide_num}", fill='gray', font=font_text)
            
            # Save image
            output_filename = f"slide_{slide_num}.{image_format}"
            output_path = os.path.join(output_dir, output_filename)
            
            if image_format.lower() == 'jpg' or image_format.lower() == 'jpeg':
                img.save(output_path, 'JPEG', quality=95)
            else:
                img.save(output_path, 'PNG')
            
            exported_slides.append({
                'filename': output_filename,
                'slide_number': slide_num,
                'size': os.path.getsize(output_path)
            })
        
        return {
            'success': True,
            'exported_count': len(exported_slides),
            'exported_slides': exported_slides,
            'output_dir': output_dir
        }
        
    except Exception as e:
        return {
            'success': False,
            'error': str(e)
        }


def export_slides_as_images_batch(input_paths, output_base_dir, image_format='png'):
    """
    Export slides from multiple PowerPoint files
    
    Args:
        input_paths: List of input PPTX paths
        output_base_dir: Base directory for outputs
        image_format: Output format
        
    Returns:
        list: Results for each presentation
    """
    results = []
    
    for input_path in input_paths:
        filename = os.path.basename(input_path)
        name, ext = os.path.splitext(filename)
        
        # Create subdirectory for this presentation's slides
        output_dir = os.path.join(output_base_dir, f"{name}_slides")
        
        result = export_slides_as_images(input_path, output_dir, image_format)
        result['ppt_filename'] = filename
        results.append(result)
    
    return results